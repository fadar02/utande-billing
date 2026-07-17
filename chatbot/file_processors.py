import os
import re
import csv
import json
import time
import hashlib
from typing import Optional, Dict, Any, List
from pathlib import Path
from datetime import datetime

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class VisionProcessor:
    """Handles image file processing and analysis."""

    SUPPORTED_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.ico'}

    def __init__(self):
        self.pil_available = PIL_AVAILABLE

    def is_image_file(self, filepath: str) -> bool:
        """Check if a file is an image based on extension."""
        ext = Path(filepath).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def analyze_image(self, filepath: str) -> str:
        """Analyze an image file. Tries PIL, falls back to file info."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        if not self.is_image_file(filepath):
            return f"Error: Not a supported image file: {filepath}"

        if self.pil_available:
            try:
                return self._analyze_with_pil(filepath)
            except Exception as e:
                return f"Error analyzing image with PIL: {e}"

        return self._get_image_info_fallback(filepath)

    def extract_text(self, filepath: str) -> str:
        """Extract text from image (OCR stub)."""
        if not self.is_image_file(filepath):
            return "Error: Not an image file"

        if self.pil_available:
            return (
                "OCR text extraction is available but requires pytesseract.\n"
                "Install with: pip install pytesseract\n"
                "Note: Tesseract-OCR must also be installed on your system.\n"
                f"File: {os.path.basename(filepath)}"
            )

        return (
            "OCR text extraction requires PIL and pytesseract.\n"
            "Install with: pip install Pillow pytesseract\n"
            f"File: {os.path.basename(filepath)}"
        )

    def _analyze_with_pil(self, filepath: str) -> str:
        """Analyze image using PIL."""
        with Image.open(filepath) as img:
            info_lines = [
                f"Image Analysis: {os.path.basename(filepath)}",
                f"  Format: {img.format}",
                f"  Mode: {img.mode}",
                f"  Size: {img.width}x{img.height} pixels",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if hasattr(img, 'info') and img.info:
                for key, val in img.info.items():
                    if isinstance(val, (str, int, float)):
                        info_lines.append(f"  {key}: {val}")

            return "\n".join(info_lines)

    def _get_image_info_fallback(self, filepath: str) -> str:
        """Get basic image info without PIL."""
        stat = os.stat(filepath)
        ext = Path(filepath).suffix.lower()

        mime_map = {
            '.png': 'PNG', '.jpg': 'JPEG', '.jpeg': 'JPEG',
            '.gif': 'GIF', '.bmp': 'BMP', '.tiff': 'TIFF',
            '.webp': 'WebP', '.ico': 'ICO'
        }

        return (
            f"Image Analysis (limited - install Pillow for full analysis): {os.path.basename(filepath)}\n"
            f"  Format: {mime_map.get(ext, 'Unknown')}\n"
            f"  File Size: {self._format_size(stat.st_size)}\n"
            f"  Last Modified: {datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"  Install Pillow for detailed analysis: pip install Pillow"
        )

    def _format_size(self, size_bytes: int) -> str:
        """Format file size to human readable."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.2f} TB"


class PDFProcessor:
    """Handles PDF file processing and analysis."""

    def __init__(self):
        self.pypdf2_available = PYPDF2_AVAILABLE

    def analyze_pdf(self, filepath: str) -> str:
        """Analyze a PDF file. Tries PyPDF2, falls back to file info."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        if not filepath.lower().endswith('.pdf'):
            return f"Error: Not a PDF file: {filepath}"

        if self.pypdf2_available:
            try:
                return self._analyze_with_pypdf2(filepath)
            except Exception as e:
                return f"Error analyzing PDF: {e}"

        return self._get_pdf_info_fallback(filepath)

    def extract_text(self, filepath: str) -> str:
        """Extract text content from a PDF."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        if not self.pypdf2_available:
            return (
                "PDF text extraction requires PyPDF2.\n"
                "Install with: pip install PyPDF2\n"
                f"File: {os.path.basename(filepath)}"
            )

        try:
            text_parts = []
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

            if not text_parts:
                return "No extractable text found in PDF (may be image-based)."

            return "\n\n".join(text_parts)
        except Exception as e:
            return f"Error extracting text from PDF: {e}"

    def summarize(self, filepath: str) -> str:
        """Summarize PDF content."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        if not self.pypdf2_available:
            return self._get_pdf_info_fallback(filepath)

        try:
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)

                all_text = ""
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        all_text += text + "\n"

                word_count = len(all_text.split())
                char_count = len(all_text)

                summary_parts = [
                    f"PDF Summary: {os.path.basename(filepath)}",
                    f"  Total Pages: {total_pages}",
                    f"  Word Count: {word_count}",
                    f"  Character Count: {char_count}",
                ]

                if all_text.strip():
                    preview = all_text.strip()[:500]
                    if len(all_text.strip()) > 500:
                        preview += "..."
                    summary_parts.append(f"\nContent Preview:\n{preview}")

                return "\n".join(summary_parts)
        except Exception as e:
            return f"Error summarizing PDF: {e}"

    def _analyze_with_pypdf2(self, filepath: str) -> str:
        """Analyze PDF using PyPDF2."""
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            info = reader.metadata
            total_pages = len(reader.pages)

            result = [
                f"PDF Analysis: {os.path.basename(filepath)}",
                f"  Pages: {total_pages}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if info:
                if info.title:
                    result.append(f"  Title: {info.title}")
                if info.author:
                    result.append(f"  Author: {info.author}")
                if info.subject:
                    result.append(f"  Subject: {info.subject}")
                if info.creator:
                    result.append(f"  Creator: {info.creator}")

            if total_pages > 0:
                first_page = reader.pages[0]
                text = first_page.extract_text()
                if text:
                    preview = text.strip()[:300]
                    result.append(f"\nFirst Page Preview:\n{preview}")

            return "\n".join(result)

    def _get_pdf_info_fallback(self, filepath: str) -> str:
        """Get basic PDF info without PyPDF2."""
        stat = os.stat(filepath)
        return (
            f"PDF Analysis (limited - install PyPDF2 for full analysis): {os.path.basename(filepath)}\n"
            f"  File Size: {self._format_size(stat.st_size)}\n"
            f"  Last Modified: {datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"  Install PyPDF2 for detailed analysis: pip install PyPDF2"
        )

    def _format_size(self, size_bytes: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.2f} TB"


class DocumentProcessor:
    """Handles Word (.docx), PowerPoint (.pptx), and Excel (.xlsx) processing."""

    SUPPORTED_EXTENSIONS = {'.docx', '.pptx', '.xlsx'}

    def __init__(self):
        self.docx_available = DOCX_AVAILABLE
        self.pptx_available = PPTX_AVAILABLE
        self.xlsx_available = OPENPYXL_AVAILABLE

    def is_supported(self, filepath: str) -> bool:
        """Check if file is a supported document type."""
        return Path(filepath).suffix.lower() in self.SUPPORTED_EXTENSIONS

    def analyze_document(self, filepath: str) -> str:
        """Analyze a document file and route to the appropriate handler."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext == '.docx':
            return self._analyze_docx(filepath)
        elif ext == '.pptx':
            return self._analyze_pptx(filepath)
        elif ext == '.xlsx':
            return self._analyze_xlsx(filepath)
        else:
            return f"Error: Unsupported document type: {ext}"

    def extract_text(self, filepath: str) -> str:
        """Extract text from a document file."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext == '.docx':
            return self._extract_docx_text(filepath)
        elif ext == '.pptx':
            return self._extract_pptx_text(filepath)
        elif ext == '.xlsx':
            return self._extract_xlsx_text(filepath)
        else:
            return f"Error: Unsupported document type: {ext}"

    def _analyze_docx(self, filepath: str) -> str:
        """Analyze a Word document."""
        if not self.docx_available:
            return (
                f"Word Analysis (limited): {os.path.basename(filepath)}\n"
                f"  File Size: {self._format_size(os.path.getsize(filepath))}\n"
                f"  Install python-docx for full analysis: pip install python-docx"
            )

        try:
            doc = DocxDocument(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables_count = len(doc.tables)
            total_words = sum(len(p.split()) for p in paragraphs)

            result = [
                f"Word Document Analysis: {os.path.basename(filepath)}",
                f"  Paragraphs: {len(paragraphs)}",
                f"  Total Words: {total_words}",
                f"  Tables: {tables_count}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if doc.core_properties.title:
                result.append(f"  Title: {doc.core_properties.title}")
            if doc.core_properties.author:
                result.append(f"  Author: {doc.core_properties.author}")

            if paragraphs:
                preview = "\n".join(paragraphs[:5])[:500]
                result.append(f"\nContent Preview:\n{preview}")

            return "\n".join(result)
        except Exception as e:
            return f"Error analyzing Word document: {e}"

    def _analyze_pptx(self, filepath: str) -> str:
        """Analyze a PowerPoint presentation."""
        if not self.pptx_available:
            return (
                f"PowerPoint Analysis (limited): {os.path.basename(filepath)}\n"
                f"  File Size: {self._format_size(os.path.getsize(filepath))}\n"
                f"  Install python-pptx for full analysis: pip install python-pptx"
            )

        try:
            prs = Presentation(filepath)
            total_slides = len(prs.slides)
            total_shapes = 0
            total_text = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    total_shapes += 1
                    if hasattr(shape, "text"):
                        total_text += shape.text + "\n"

            word_count = len(total_text.split())

            result = [
                f"PowerPoint Analysis: {os.path.basename(filepath)}",
                f"  Slides: {total_slides}",
                f"  Total Shapes: {total_shapes}",
                f"  Word Count: {word_count}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if total_text.strip():
                preview = total_text.strip()[:500]
                result.append(f"\nContent Preview:\n{preview}")

            return "\n".join(result)
        except Exception as e:
            return f"Error analyzing PowerPoint: {e}"

    def _analyze_xlsx(self, filepath: str) -> str:
        """Analyze an Excel spreadsheet."""
        if not self.xlsx_available:
            return (
                f"Excel Analysis (limited): {os.path.basename(filepath)}\n"
                f"  File Size: {self._format_size(os.path.getsize(filepath))}\n"
                f"  Install openpyxl for full analysis: pip install openpyxl"
            )

        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            total_rows = 0
            total_cols = 0

            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                total_rows += ws.max_row or 0
                total_cols += ws.max_column or 0

            wb.close()

            result = [
                f"Excel Analysis: {os.path.basename(filepath)}",
                f"  Sheets: {len(sheet_names)}",
                f"  Sheet Names: {', '.join(sheet_names)}",
                f"  Total Rows: {total_rows}",
                f"  Total Columns: {total_cols}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            return "\n".join(result)
        except Exception as e:
            return f"Error analyzing Excel file: {e}"

    def _extract_docx_text(self, filepath: str) -> str:
        """Extract text from Word document."""
        if not self.docx_available:
            return "Text extraction requires python-docx. Install with: pip install python-docx"

        try:
            doc = DocxDocument(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else "No text content found."
        except Exception as e:
            return f"Error extracting text from Word document: {e}"

    def _extract_pptx_text(self, filepath: str) -> str:
        """Extract text from PowerPoint."""
        if not self.pptx_available:
            return "Text extraction requires python-pptx. Install with: pip install python-pptx"

        try:
            prs = Presentation(filepath)
            all_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    all_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
            return "\n\n".join(all_text) if all_text else "No text content found."
        except Exception as e:
            return f"Error extracting text from PowerPoint: {e}"

    def _extract_xlsx_text(self, filepath: str) -> str:
        """Extract text from Excel."""
        if not self.xlsx_available:
            return "Text extraction requires openpyxl. Install with: pip install openpyxl"

        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            all_text = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_text = []
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(cell) for cell in row if cell is not None]
                    if row_vals:
                        rows_text.append(", ".join(row_vals))
                if rows_text:
                    all_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text))

            wb.close()
            return "\n\n".join(all_text) if all_text else "No text content found."
        except Exception as e:
            return f"Error extracting text from Excel: {e}"

    def _format_size(self, size_bytes: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.2f} TB"


class DataProcessor:
    """Handles CSV, JSON, and Excel data for analysis."""

    def __init__(self):
        self.xlsx_available = OPENPYXL_AVAILABLE

    def analyze_data(self, filepath: str) -> str:
        """Analyze a data file and provide summary statistics."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext == '.csv':
            return self._analyze_csv(filepath)
        elif ext == '.json':
            return self._analyze_json(filepath)
        elif ext == '.xlsx' or ext == '.xls':
            return self._analyze_excel_data(filepath)
        else:
            return f"Error: Unsupported data file type: {ext}"

    def create_chart_description(self, filepath: str) -> str:
        """Create a text description suitable for chart generation."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext == '.csv':
            return self._chart_description_csv(filepath)
        elif ext == '.json':
            return self._chart_description_json(filepath)
        elif ext == '.xlsx' or ext == '.xls':
            return self._chart_description_excel(filepath)
        else:
            return f"Error: Unsupported file type for charting: {ext}"

    def summarize_data(self, filepath: str) -> str:
        """Provide a comprehensive data summary."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext == '.csv':
            return self._summarize_csv(filepath)
        elif ext == '.json':
            return self._summarize_json(filepath)
        elif ext == '.xlsx' or ext == '.xls':
            return self._summarize_excel(filepath)
        else:
            return f"Error: Unsupported file type: {ext}"

    def _analyze_csv(self, filepath: str) -> str:
        """Analyze a CSV file."""
        try:
            with open(filepath, 'r', newline='', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                rows = list(reader)

            if not headers:
                return f"CSV file appears empty: {os.path.basename(filepath)}"

            result = [
                f"CSV Analysis: {os.path.basename(filepath)}",
                f"  Columns: {len(headers)}",
                f"  Rows: {len(rows)}",
                f"  Headers: {', '.join(headers)}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if rows:
                result.append(f"\n  Sample Data (first 3 rows):")
                for i, row in enumerate(rows[:3]):
                    row_str = ", ".join(row[:len(headers)])
                    result.append(f"    Row {i + 1}: {row_str}")

            return "\n".join(result)
        except Exception as e:
            return f"Error analyzing CSV: {e}"

    def _analyze_json(self, filepath: str) -> str:
        """Analyze a JSON file."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                data = json.load(f)

            result = [
                f"JSON Analysis: {os.path.basename(filepath)}",
                f"  Type: {type(data).__name__}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if isinstance(data, list):
                result.append(f"  Array Length: {len(data)}")
                if data:
                    first = data[0]
                    if isinstance(first, dict):
                        result.append(f"  Fields: {', '.join(first.keys())}")
            elif isinstance(data, dict):
                result.append(f"  Top-level Keys: {', '.join(data.keys())}")
                for key, val in list(data.items())[:3]:
                    val_type = type(val).__name__
                    if isinstance(val, (list, dict)):
                        result.append(f"  {key}: {val_type} with {len(val)} items")
                    else:
                        result.append(f"  {key}: {val_type}")

            return "\n".join(result)
        except json.JSONDecodeError:
            return f"Error: Invalid JSON in {os.path.basename(filepath)}"
        except Exception as e:
            return f"Error analyzing JSON: {e}"

    def _analyze_excel_data(self, filepath: str) -> str:
        """Analyze Excel file as data."""
        if not self.xlsx_available:
            return (
                f"Excel Data Analysis (limited): {os.path.basename(filepath)}\n"
                f"  File Size: {self._format_size(os.path.getsize(filepath))}\n"
                f"  Install openpyxl for full analysis: pip install openpyxl"
            )

        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            result = [
                f"Excel Data Analysis: {os.path.basename(filepath)}",
                f"  Sheets: {len(wb.sheetnames)}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if rows:
                    headers = [str(cell) for cell in rows[0] if cell is not None]
                    data_rows = rows[1:]
                    result.append(f"\n  Sheet: {sheet_name}")
                    result.append(f"    Columns: {', '.join(headers)}")
                    result.append(f"    Data Rows: {len(data_rows)}")

            wb.close()
            return "\n".join(result)
        except Exception as e:
            return f"Error analyzing Excel data: {e}"

    def _chart_description_csv(self, filepath: str) -> str:
        """Generate chart description from CSV data."""
        try:
            with open(filepath, 'r', newline='', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                rows = list(reader)

            if not headers or not rows:
                return "Insufficient data for chart generation."

            numeric_cols = []
            for i, header in enumerate(headers):
                numeric_count = 0
                for row in rows[:min(100, len(rows))]:
                    if i < len(row):
                        try:
                            float(row[i])
                            numeric_count += 1
                        except (ValueError, TypeError):
                            pass
                if numeric_count > len(rows[:min(100, len(rows))]) * 0.5:
                    numeric_cols.append((i, header))

            result = [
                f"Chart Data Description: {os.path.basename(filepath)}",
                f"  Available for Charting:",
            ]

            if numeric_cols:
                for idx, name in numeric_cols:
                    values = []
                    for row in rows:
                        if idx < len(row):
                            try:
                                values.append(float(row[idx]))
                            except (ValueError, TypeError):
                                pass
                    if values:
                        result.append(f"    Column '{name}': min={min(values):.2f}, max={max(values):.2f}, avg={sum(values)/len(values):.2f}")
            else:
                result.append("    No numeric columns detected for charting.")

            label_col = headers[0] if headers else None
            if label_col:
                unique_labels = set(row[0] for row in rows[:50] if row)
                result.append(f"  Suggested Label Column: '{label_col}' ({len(unique_labels)} unique values)")

            return "\n".join(result)
        except Exception as e:
            return f"Error generating chart description: {e}"

    def _chart_description_json(self, filepath: str) -> str:
        """Generate chart description from JSON data."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                data = json.load(f)

            if isinstance(data, list) and data and isinstance(data[0], dict):
                headers = list(data[0].keys())
                result = [
                    f"Chart Data Description: {os.path.basename(filepath)}",
                    f"  Data Points: {len(data)}",
                    f"  Fields: {', '.join(headers)}",
                ]
                for key in headers:
                    values = [item[key] for item in data[:50] if key in item and isinstance(item[key], (int, float))]
                    if values:
                        result.append(f"    '{key}': min={min(values):.2f}, max={max(values):.2f}, avg={sum(values)/len(values):.2f}")
                return "\n".join(result)

            return f"JSON structure not ideal for charting. Type: {type(data).__name__}"
        except Exception as e:
            return f"Error generating chart description: {e}"

    def _chart_description_excel(self, filepath: str) -> str:
        """Generate chart description from Excel data."""
        if not self.xlsx_available:
            return "Chart description requires openpyxl. Install with: pip install openpyxl"
        return f"Excel chart description: use the CSV/JSON chart methods after exporting."

    def _summarize_csv(self, filepath: str) -> str:
        """Provide comprehensive CSV summary."""
        try:
            with open(filepath, 'r', newline='', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                rows = list(reader)

            if not headers:
                return "CSV file is empty."

            total_cells = len(headers) * len(rows)
            empty_cells = sum(1 for row in rows for cell in row if not cell.strip())

            result = [
                f"CSV Summary: {os.path.basename(filepath)}",
                f"  Rows: {len(rows)}",
                f"  Columns: {len(headers)}",
                f"  Total Cells: {total_cells}",
                f"  Empty Cells: {empty_cells} ({empty_cells/total_cells*100:.1f}%)" if total_cells else "  Empty Cells: 0",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            result.append(f"\n  Column Details:")
            for i, header in enumerate(headers):
                values = [row[i] for row in rows if i < len(row)]
                non_empty = [v for v in values if v.strip()]
                unique_vals = len(set(non_empty))
                numeric_vals = sum(1 for v in non_empty if self._is_numeric(v))
                result.append(f"    '{header}': {len(non_empty)}/{len(values)} filled, {unique_vals} unique, {numeric_vals} numeric")

            return "\n".join(result)
        except Exception as e:
            return f"Error summarizing CSV: {e}"

    def _summarize_json(self, filepath: str) -> str:
        """Provide comprehensive JSON summary."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                data = json.load(f)

            result = [
                f"JSON Summary: {os.path.basename(filepath)}",
                f"  Root Type: {type(data).__name__}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
            ]

            if isinstance(data, list):
                result.append(f"  Array Length: {len(data)}")
                types = {}
                for item in data:
                    t = type(item).__name__
                    types[t] = types.get(t, 0) + 1
                result.append(f"  Element Types: {types}")

                if data and isinstance(data[0], dict):
                    all_keys = {}
                    for item in data:
                        if isinstance(item, dict):
                            for key in item:
                                all_keys[key] = all_keys.get(key, 0) + 1
                    result.append(f"  Key Frequency: {all_keys}")

            elif isinstance(data, dict):
                result.append(f"  Keys: {len(data)}")
                for key, val in data.items():
                    result.append(f"    '{key}': {type(val).__name__}")

            return "\n".join(result)
        except Exception as e:
            return f"Error summarizing JSON: {e}"

    def _summarize_excel(self, filepath: str) -> str:
        """Provide comprehensive Excel summary."""
        if not self.xlsx_available:
            return "Excel summary requires openpyxl. Install with: pip install openpyxl"

        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            result = [
                f"Excel Summary: {os.path.basename(filepath)}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
                f"  Sheets: {len(wb.sheetnames)}",
            ]

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if rows:
                    headers = [str(cell) for cell in rows[0] if cell is not None]
                    result.append(f"\n  Sheet: '{sheet_name}'")
                    result.append(f"    Headers: {', '.join(headers)}")
                    result.append(f"    Data Rows: {len(rows) - 1}")

            wb.close()
            return "\n".join(result)
        except Exception as e:
            return f"Error summarizing Excel: {e}"

    def _is_numeric(self, value: str) -> bool:
        """Check if a string value is numeric."""
        try:
            float(value)
            return True
        except (ValueError, TypeError):
            return False

    def _format_size(self, size_bytes: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.2f} TB"


class FileProcessor:
    """Main file processor that orchestrates all specialized processors."""

    FILE_REQUEST_PATTERNS = [
        r'analyze\s+(?:this\s+)?(?:file|document)',
        r'process\s+(?:this\s+)?(?:file|document)',
        r'read\s+(?:this\s+)?(?:file|document)',
        r'open\s+(?:this\s+)?(?:file|document)',
        r'what(?:\'s|\s+is)\s+in\s+(?:this\s+)?(?:file|document)',
        r'show\s+(?:me\s+)?(?:the\s+)?(?:file|document)',
        r'extract\s+(?:text|data)\s+from',
        r'summarize\s+(?:this\s+)?(?:file|document|data)',
        r'describe\s+(?:this\s+)?(?:file|document|image|pdf)',
        r'\.pdf\b',
        r'\.docx\b',
        r'\.xlsx\b',
        r'\.pptx\b',
        r'\.csv\b',
        r'\.json\b',
        r'\.png\b',
        r'\.jpg\b',
        r'\.jpeg\b',
        r'\.gif\b',
    ]

    def __init__(self):
        self.vision = VisionProcessor()
        self.pdf = PDFProcessor()
        self.document = DocumentProcessor()
        self.data = DataProcessor()
        self._compiled_patterns = [re.compile(p, re.IGNORECASE) for p in self.FILE_REQUEST_PATTERNS]

    def is_file_request(self, text: str) -> bool:
        """Detect if user input is requesting file processing."""
        if not text or not text.strip():
            return False

        text_lower = text.lower().strip()

        for pattern in self._compiled_patterns:
            if pattern.search(text_lower):
                return True

        words = text_lower.split()
        for word in words:
            clean_word = word.strip('.,!?()[]{}"\'')
            ext = Path(clean_word).suffix.lower()
            if ext in {'.pdf', '.docx', '.xlsx', '.pptx', '.csv', '.json',
                       '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.txt', '.py'}:
                return True

        if os.path.exists(text.strip()):
            return True

        return False

    def handle(self, text: str) -> str:
        """Handle a file processing request from user text."""
        text = text.strip()

        filepath = self._extract_filepath(text)
        if filepath:
            return self.process_file(filepath)

        return (
            "I can process various file types! Please provide a file path.\n\n"
            "Supported formats:\n"
            "  Images: PNG, JPG, GIF, BMP, WEBP\n"
            "  Documents: PDF, DOCX, PPTX, XLSX\n"
            "  Data: CSV, JSON\n\n"
            "Example: 'analyze /path/to/file.pdf'"
        )

    def get_file_info(self, filepath: str) -> str:
        """Get basic information about a file."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        stat = os.stat(filepath)
        ext = Path(filepath).suffix.lower()
        name = os.path.basename(filepath)

        result = [
            f"File Information: {name}",
            f"  Full Path: {filepath}",
            f"  Extension: {ext or 'None'}",
            f"  Size: {self._format_size(stat.st_size)}",
            f"  Created: {datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S')}",
            f"  Modified: {datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}",
            f"  Is File: {os.path.isfile(filepath)}",
            f"  Is Directory: {os.path.isdir(filepath)}",
            f"  Readable: {os.access(filepath, os.R_OK)}",
        ]

        return "\n".join(result)

    def process_file(self, filepath: str) -> str:
        """Process a file by detecting its type and routing to the appropriate processor."""
        if not os.path.exists(filepath):
            return f"Error: File not found: {filepath}"

        if not os.path.isfile(filepath):
            if os.path.isdir(filepath):
                return self._list_directory(filepath)
            return f"Error: Path is not a file: {filepath}"

        ext = Path(filepath).suffix.lower()

        if ext in VisionProcessor.SUPPORTED_EXTENSIONS:
            return self.vision.analyze_image(filepath)

        if ext == '.pdf':
            return self.pdf.analyze_pdf(filepath)

        if ext in {'.docx', '.pptx', '.xlsx'}:
            return self.document.analyze_document(filepath)

        if ext == '.csv':
            return self.data.analyze_data(filepath)

        if ext == '.json':
            return self.data.analyze_data(filepath)

        if ext == '.txt':
            return self._read_text_file(filepath)

        if ext == '.py':
            return self._read_code_file(filepath, 'Python')

        if ext in {'.js', '.ts'}:
            return self._read_code_file(filepath, 'JavaScript/TypeScript')

        return self._read_generic_file(filepath)

    def _extract_filepath(self, text: str) -> Optional[str]:
        """Extract a file path from user text."""
        quoted_match = re.search(r'["\']([^"\']+)["\']', text)
        if quoted_match:
            path = quoted_match.group(1)
            if os.path.exists(path):
                return path

        path_match = re.search(r'(/[\w./\-_~]+|[\w./\\\-_~]+\.\w{1,5})', text)
        if path_match:
            path = path_match.group(1)
            if os.path.exists(path):
                return path

        words = text.split()
        for word in words:
            clean = word.strip('.,!?()[]{}"\'')
            if os.path.exists(clean):
                return clean
            if os.path.exists(os.path.join('.', clean)):
                return os.path.join('.', clean)

        return None

    def _read_text_file(self, filepath: str) -> str:
        """Read and return contents of a text file."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()

            lines = content.split('\n')
            result = [
                f"Text File: {os.path.basename(filepath)}",
                f"  Lines: {len(lines)}",
                f"  Characters: {len(content)}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
                f"\nContent:\n{content[:5000]}",
            ]

            if len(content) > 5000:
                result.append(f"\n... [truncated, showing first 5000 of {len(content)} characters]")

            return "\n".join(result)
        except Exception as e:
            return f"Error reading text file: {e}"

    def _read_code_file(self, filepath: str, language: str) -> str:
        """Read and display a code file."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()

            lines = content.split('\n')
            result = [
                f"{language} File: {os.path.basename(filepath)}",
                f"  Lines: {len(lines)}",
                f"  File Size: {self._format_size(os.path.getsize(filepath))}",
                f"\nCode Preview:\n{content[:3000]}",
            ]

            if len(content) > 3000:
                result.append(f"\n... [truncated, showing first 3000 of {len(content)} characters]")

            return "\n".join(result)
        except Exception as e:
            return f"Error reading code file: {e}"

    def _read_generic_file(self, filepath: str) -> str:
        """Attempt to read a generic file."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(10000)

            stat = os.stat(filepath)
            return (
                f"File: {os.path.basename(filepath)}\n"
                f"  Size: {self._format_size(stat.st_size)}\n"
                f"  Content Preview (first 10000 chars):\n{content}"
            )
        except UnicodeDecodeError:
            stat = os.stat(filepath)
            return (
                f"Binary File: {os.path.basename(filepath)}\n"
                f"  Size: {self._format_size(stat.st_size)}\n"
                f"  This appears to be a binary file and cannot be displayed as text."
            )
        except Exception as e:
            return f"Error reading file: {e}"

    def _list_directory(self, dirpath: str) -> str:
        """List contents of a directory."""
        try:
            entries = os.listdir(dirpath)
            result = [f"Directory: {dirpath}", f"  Total Entries: {len(entries)}", ""]

            dirs = sorted([e for e in entries if os.path.isdir(os.path.join(dirpath, e))])
            files = sorted([e for e in entries if os.path.isfile(os.path.join(dirpath, e))])

            if dirs:
                result.append("  Directories:")
                for d in dirs[:20]:
                    result.append(f"    {d}/")
                if len(dirs) > 20:
                    result.append(f"    ... and {len(dirs) - 20} more directories")

            if files:
                result.append("\n  Files:")
                for f in files[:20]:
                    size = os.path.getsize(os.path.join(dirpath, f))
                    result.append(f"    {f} ({self._format_size(size)})")
                if len(files) > 20:
                    result.append(f"    ... and {len(files) - 20} more files")

            return "\n".join(result)
        except Exception as e:
            return f"Error listing directory: {e}"

    def _format_size(self, size_bytes: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.2f} TB"
